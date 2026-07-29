import fitz
from pptx import Presentation
import tempfile
import sys
import os

pdf = sys.argv[1]
pptx = sys.argv[2]

doc = fitz.open(pdf)

prs = Presentation()

# Remove o slide inicial vazio
while len(prs.slides._sldIdLst) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

for i, page in enumerate(doc):

    pix = page.get_pixmap(dpi=200)

    temp_path = os.path.join(
        tempfile.gettempdir(),
        f"pdfprime_page_{i}.png"
    )

    pix.save(temp_path)

    slide = prs.slides.add_slide(
        prs.slide_layouts[6]
    )

    slide.shapes.add_picture(
        temp_path,
        0,
        0,
        width=prs.slide_width,
        height=prs.slide_height
    )

    if os.path.exists(temp_path):
        os.remove(temp_path)

prs.save(pptx)

doc.close()

print("Conversão concluída.")